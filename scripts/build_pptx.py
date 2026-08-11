"""manual-draft.md → 매뉴얼 PPTX 생성기 (전용 pptx skill 이 없는 환경의 표준 폴백).

output-formats.md 의 pptx 명세를 내장한다:
  표지 → CONTENTS(장·절 목록 + 슬라이드 번호) → 장 간지 → 화면 슬라이드(1절=1화면),
  설명 6개 초과 시 (1)/(2) 분할, 마크다운 서식은 실제 서식으로 변환(기호 잔존 금지),
  placeholder 는 이미지 프레임과 동일 크기의 회색 상자 + 캡션, 페이지 번호는 PPT 슬라이드 순번.
  이미지 없는 짧은 절(시스템 개요·권한 체계·접속 환경 등)이 같은 장에 연속되면
  한 슬라이드에 병합한다(절 제목을 소제목으로 스택) — 절마다 거의 빈 장이 생기는 낭비 방지.
  제목만 있고 본문이 없는 부모 절(3단 번호의 그룹 제목)은 슬라이드를 만들지 않고
  CONTENTS 번호만 첫 하위 절 슬라이드로 위임한다 — 헤더만 있는 빈 장표 방지.
  --orientation 으로 세로(A4, 기본)/가로(16:9) 를 선택할 수 있다. 세로형은 모든
  캡처를 상(이미지)/하(설명)으로 배치하고 CONTENTS 를 1컬럼으로 구성한다.
  --theme 으로 색 테마(navy 기본/forest/charcoal)를 선택할 수 있다 — 표지·간지·
  헤더의 배경/포인트 팔레트만 바뀌고 본문 텍스트·주의(※) 색은 공통이다.

사용 예:
  python build_pptx.py --draft manual-work/manual-draft.md \
      --screenshots manual-work/screenshots --out 관리자매뉴얼_시스템명_20260711.pptx \
      [--orientation portrait]

종료 코드: 0 성공 / 1 파싱·생성 오류 / 2 python-pptx 미설치
"""

import argparse
import math
import os
import re
import sys

if hasattr(sys.stdout, "reconfigure"):
    sys.stdout.reconfigure(encoding="utf-8")
    sys.stderr.reconfigure(encoding="utf-8")

sys.path.insert(0, os.path.dirname(os.path.abspath(__file__)))
from draft_parser import (parse_draft, parse_inline, parse_meta, plain, image_size,
                          resolve_image, text_lines, tile_tall_image, CIRCLED,
                          table_height_est, tables_height, TABLE_PAD,
                          PORT_BODY_W, PORT_TEXT_BOTTOM, PORT_IMG_Y)


def fail(msg, code=1):
    print(f"[build_pptx] 오류: {msg}", file=sys.stderr)
    sys.exit(code)


try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.oxml.ns import qn
except ImportError:
    fail("python-pptx 가 설치되어 있지 않습니다. 설치 후 재시도하세요:\n  pip install python-pptx", code=2)

TEXT = RGBColor(0x26, 0x26, 0x26)
MUTED = RGBColor(0x8A, 0x8F, 0x98)
NOTE = RGBColor(0xC0, 0x39, 0x2B)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
PH_BG = RGBColor(0xEC, 0xEE, 0xF1)
PH_TX = RGBColor(0x6B, 0x72, 0x80)
FONT = "맑은 고딕"

# 테마 팔레트 — dark(표지·간지·헤더 배경), accent(포인트·접근 경로·액센트 바),
# soft/softer(어두운 배경 위 보조 텍스트), head_num(헤더의 절 번호).
# 본문 텍스트(TEXT)·주의(NOTE)·placeholder 색은 의미색이라 테마와 무관하게 공통이다.
THEMES = {
    "navy":     {"dark": "1F2A44", "accent": "2E5BFF",
                 "soft": "B9C4DE", "softer": "9AA5C0", "head_num": "7E96E8"},
    "forest":   {"dark": "1F3D2F", "accent": "0FA47A",
                 "soft": "BCD8CB", "softer": "98BFAF", "head_num": "6FC9A8"},
    "charcoal": {"dark": "2B2F36", "accent": "E8590C",
                 "soft": "C9CED6", "softer": "A6ADB8", "head_num": "F08C4B"},
}


def _rgb(hex6):
    return RGBColor(int(hex6[0:2], 16), int(hex6[2:4], 16), int(hex6[4:6], 16))


def apply_theme(name):
    """표지·간지·헤더의 배경/포인트 팔레트를 모듈 전역에 적용한다.

    ACCENT 는 흰 배경 위(본문 '접근 경로'·CONTENTS 바), ACCENT_ON_DARK 는 어두운
    배경 위(표지 바·간지 장 번호) 용도다 — 기본 3종은 두 값이 같고, 템플릿 추출
    경로에서만 갈린다(밝은 accent 는 흰 배경에서 안 읽히므로 본문용만 어둡게)."""
    global THEME, DARK, ACCENT, ACCENT_ON_DARK, DARK_SOFT, DARK_SOFTER, HEAD_NUM
    t = THEMES[name]
    THEME = name
    DARK, ACCENT = _rgb(t["dark"]), _rgb(t["accent"])
    ACCENT_ON_DARK = ACCENT
    DARK_SOFT, DARK_SOFTER, HEAD_NUM = _rgb(t["soft"]), _rgb(t["softer"]), _rgb(t["head_num"])


def _hls_hex(h, l, s):
    import colorsys
    r, g, b = colorsys.hls_to_rgb(h, l, s)
    return f"{round(r * 255):02X}{round(g * 255):02X}{round(b * 255):02X}"


def _hex_hls(hex6):
    import colorsys
    return colorsys.rgb_to_hls(int(hex6[0:2], 16) / 255, int(hex6[2:4], 16) / 255,
                               int(hex6[4:6], 16) / 255)


def _contrast_on_white(hex6):
    """흰 배경 대비 명도비(WCAG). 4.5 이상이면 본문 텍스트로 읽을 만하다."""
    def lin(c):
        c = int(c, 16) / 255
        return c / 12.92 if c <= 0.03928 else ((c + 0.055) / 1.055) ** 2.4
    lum = 0.2126 * lin(hex6[0:2]) + 0.7152 * lin(hex6[2:4]) + 0.0722 * lin(hex6[4:6])
    return 1.05 / (lum + 0.05)


def _darken_for_white_bg(hex6, target=4.5):
    """흰 배경 위 본문 색으로 쓸 수 있게 명도를 낮춘다 — 노랑처럼 색상 자체가 밝은
    계열은 고정 명도로는 대비가 안 나오므로 목표 명도비에 도달할 때까지 단계적으로
    어둡게 한다(색상·채도는 유지해 템플릿 색감을 잃지 않는다)."""
    h, l, s = _hex_hls(hex6)
    out = hex6
    while l > 0.12 and _contrast_on_white(out) < target:
        l -= 0.04
        out = _hls_hex(h, l, max(s, 0.35))
    return out


def theme_from_template(path):
    """참고 템플릿 pptx 의 색 테마(theme1.xml clrScheme)에서 dark·accent 를 추출해
    커스텀 팔레트를 파생·적용한다. **레이아웃·규격은 번들 그대로** — '템플릿 참고'는
    슬라이드 서식 복제가 아니라 색 스타일 추출이다. 서식 복제는 레이아웃 상속·
    placeholder 기하 등에서 디자인이 깨지기 쉬워 명시 요청 시의 예외 경로로만 둔다.

    성공 시 (dark_hex, accent_hex) 반환, 실패 시 None(호출부가 기본 테마 유지)."""
    import zipfile
    try:
        with zipfile.ZipFile(path) as z:
            xml = z.read("ppt/theme/theme1.xml").decode("utf-8", "ignore")
    except Exception:
        return None
    m = re.search(r"<a:clrScheme.*?</a:clrScheme>", xml, re.S)
    if not m:
        return None
    scheme = m.group(0)

    def pick(tag):
        mm = re.search(rf"<a:{tag}>.*?(?:srgbClr val=\"([0-9A-Fa-f]{{6}})\""
                       rf"|sysClr[^>]*lastClr=\"([0-9A-Fa-f]{{6}})\")", scheme, re.S)
        return (mm.group(1) or mm.group(2)).upper() if mm else None

    dark, accent = pick("dk2"), pick("accent1")
    if not dark or not accent:
        return None
    # 무채색(회색조) 템플릿: hue 가 0 으로 떨어져 원본에 없는 색조(붉은빛)를 발명하게 되므로
    # 파생을 포기하고 기본 테마를 유지한다. 둘 중 **하나라도** 무채색이면 그 축의 파생이
    # 왜곡되므로 OR 로 판정한다(AND 로 두면 유채색 dark + 무채색 accent 조합이 빠져나간다).
    _, _, ds0 = _hex_hls(dark)
    _, _, as0 = _hex_hls(accent)
    if ds0 < 0.08 or as0 < 0.08:
        return None
    # DARK 는 표지·간지의 '배경'이면서 CONTENTS·헤딩의 '흰 배경 위 텍스트'로도 쓰인다 —
    # 밝은 dk2(예: Office Gold)면 배경 위 흰 글씨도, 흰 배경 위 제목도 모두 안 읽힌다.
    # 대비 기준을 만족할 때까지 어둡게 내린다(색상·채도 유지).
    dh, dl, ds = _hex_hls(dark)
    if dl > 0.5 or _contrast_on_white(dark) < 4.5:
        dark = _darken_for_white_bg(dark)
        dh, dl, ds = _hex_hls(dark)
    ah, al, as_ = _hex_hls(accent)
    # ACCENT 는 흰 배경 위 본문 색으로도 쓰인다('접근 경로'·CONTENTS 액센트 바) —
    # 밝은 accent(노랑·라이트그린 등)를 그대로 쓰면 대비가 1.4:1 수준이라 읽히지 않는다.
    # 본문용만 대비 기준까지 어둡게 하고, 표지·간지의 원색은 ACCENT_ON_DARK 로 보존한다.
    accent_text = _darken_for_white_bg(accent)
    # 반대 방향 가드: 어두운 accent(예: Blue Darker 50%)는 어두운 표지·간지 위에서
    # 보이지 않으므로, 배경(dark)보다 충분히 밝지 않으면 명도를 올려 쓴다
    accent_dark_bg = accent
    ad_h, ad_l, ad_s = _hex_hls(accent)
    if ad_l < dl + 0.18:
        accent_dark_bg = _hls_hex(ad_h, min(0.72, dl + 0.30), max(0.35, ad_s))
    # 보조 팔레트 파생 — 기본 3종 테마의 명도·채도 관계를 따른다
    global THEME, DARK, ACCENT, ACCENT_ON_DARK, DARK_SOFT, DARK_SOFTER, HEAD_NUM
    THEME = "template"
    DARK, ACCENT, ACCENT_ON_DARK = _rgb(dark), _rgb(accent_text), _rgb(accent_dark_bg)
    DARK_SOFT = _rgb(_hls_hex(dh, 0.80, min(0.45, max(0.12, ds))))
    DARK_SOFTER = _rgb(_hls_hex(dh, 0.68, min(0.30, max(0.10, ds * 0.7))))
    HEAD_NUM = _rgb(_hls_hex(ah, 0.70, min(0.75, max(0.25, as_))))
    return dark, accent_text


apply_theme("navy")

MAX_ITEMS = 6                 # 슬라이드당 설명 항목 상한 (초과 시 (1)/(2) 분할)
CAP_H = Inches(0.32)          # 캡션 줄 높이

PORTRAIT = False


def apply_orientation(portrait):
    """슬라이드 방향별 레이아웃 프로파일을 모듈 전역에 적용한다.

    세로(기본, A4 7.5x10.833in): 폭이 좁아 모든 캡처를 상(이미지)/하(설명)으로
    배치하고, CONTENTS 는 1컬럼으로 구성한다. 줄당 문자 수(EA)·설명 예산도 폭에
    비례해 조정한다.
    가로(16:9 13.333x7.5in): 세로 비율 캡처는 좌(이미지)/우(설명),
    가로 비율 캡처는 상(이미지)/하(설명)으로 배치한다.
    """
    global PORTRAIT, SLIDE_W, SLIDE_H, BODY_X, BODY_W, SIDE_X, SIDE_W, TEXT_BOTTOM
    global IMG_Y, V_FRAME_W, V_FRAME_H, H_FRAME_W, H_FRAME_H, PORT_IMG_MAX_H, IMG_MIN_H
    global SIDE_EA, WIDE_EA, INTRO_EA, PLAIN_LINES, TALL_RATIO_MIN
    global COMBINE_BUDGET, TOC_COLS, TOC_PER_COL, TOC_COL_W
    PORTRAIT = portrait
    # 폭 균일 임계: 이미지 비율(가로/세로)이 이 값 미만이면 전폭 렌더 시 높이 상한에
    # 걸려 폭이 줄어든다(=매뉴얼 내 다른 캡처와 크기 불일치). 그 아래는 타일 분할한다.
    # 세로형 = BODY_W/PORT_IMG_MAX_H. 가로형은 폭 축소 압력이 낮아 미사용(None).
    TALL_RATIO_MIN = round(6.4 / 5.7, 3) if portrait else None
    if portrait:
        SLIDE_W, SLIDE_H = Inches(7.5), Inches(10.833)
        # 세로형 기하는 draft_parser 가 단일 출처 — 원고 린터(validate_draft)가 pptx
        # 의존 없이 같은 값으로 표 높이를 예측하므로, 여기서 직접 바꾸면 린트가 어긋난다
        BODY_X, BODY_W = Inches(0.55), Inches(PORT_BODY_W)   # 본문 전폭
        SIDE_X, SIDE_W = None, None                      # 좌/우 분할 미사용
        TEXT_BOTTOM = Inches(PORT_TEXT_BOTTOM)           # 설명 프레임 하한 (페이지 번호 위)
        IMG_Y = Inches(PORT_IMG_Y)                       # 이미지 프레임 상단 고정
        # 세로형 이미지는 본문 폭을 가득 채우고 높이는 캡처 비율을 따른다(가변).
        # 프레임 h 는 placeholder(비율을 모름) 렌더에만 쓰는 기본값이다.
        V_FRAME_W, V_FRAME_H = Inches(6.4), Inches(4.32)
        H_FRAME_W, H_FRAME_H = Inches(6.4), Inches(4.32)
        PORT_IMG_MAX_H = Inches(5.7)                     # 세로로 긴 캡처의 높이 상한(설명 공간 확보)
        IMG_MIN_H = 3.6                                  # 설명 수용 위한 이미지 축소 하한(in)
        SIDE_EA, WIDE_EA, INTRO_EA = 33, 37, 35          # 전각 기준 줄당 문자 수(폭 비례)
        PLAIN_LINES = 30                                 # 시각 요소 없는 절의 설명 줄 예산
        COMBINE_BUDGET = 8.6                             # 개요 병합 본문 가용 높이(in)
        TOC_COLS, TOC_PER_COL, TOC_COL_W = 1, 26, Inches(6.4)
    else:
        SLIDE_W, SLIDE_H = Inches(13.333), Inches(7.5)
        BODY_X, BODY_W = Inches(0.55), Inches(12.2)
        SIDE_X, SIDE_W = Inches(7.1), Inches(5.7)        # 우측 설명 컬럼
        TEXT_BOTTOM = Inches(7.0)
        IMG_Y = Inches(2.25)
        V_FRAME_W, V_FRAME_H = Inches(6.3), Inches(4.35)  # 세로 비율 캡처(좌측) 프레임
        H_FRAME_W, H_FRAME_H = Inches(9.4), Inches(3.2)   # 가로 비율 캡처(상단) 프레임
        PORT_IMG_MAX_H = None                             # 가로형 미사용
        IMG_MIN_H = 2.3                                   # 설명 수용 위한 이미지 축소 하한(in)
        SIDE_EA, WIDE_EA, INTRO_EA = 33, 72, 68
        PLAIN_LINES = 16
        COMBINE_BUDGET = 5.6
        TOC_COLS, TOC_PER_COL, TOC_COL_W = 2, 17, Inches(5.9)


apply_orientation(False)


def _set_font(run, size, bold=False, color=TEXT):
    run.font.name = FONT
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    rPr = run._r.get_or_add_rPr()
    ea = rPr.find(qn("a:ea"))
    if ea is None:
        ea = rPr.makeelement(qn("a:ea"), {})
        rPr.append(ea)
    ea.set("typeface", FONT)


def add_text(slide, x, y, w, h, wrap=True, anchor=MSO_ANCHOR.TOP):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = wrap
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = Inches(0.02)
    tf.margin_top = tf.margin_bottom = Inches(0.01)
    return tf


def set_para(p, segments, size, color=TEXT, bold=False, align=PP_ALIGN.LEFT, space_after=4):
    """segments: 문자열 또는 (text, bold) 목록."""
    p.alignment = align
    p.space_after = Pt(space_after)
    if isinstance(segments, str):
        segments = [(segments, bold)]
    for text, seg_bold in segments:
        r = p.add_run()
        r.text = text
        _set_font(r, size, bold=bold or seg_bold, color=color)
    return p


def add_para(tf, *args, **kwargs):
    p = tf.paragraphs[0] if not tf.paragraphs[0].runs else tf.add_paragraph()
    return set_para(p, *args, **kwargs)


def add_rect(slide, x, y, w, h, fill, line=None):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        shp.line.width = Pt(0.75)
    shp.shadow.inherit = False
    return shp


def image_ratio(path):
    """가로/세로 비율 — 치수를 모르는 포맷은 배치 판정용 기본값 1.33 (렌더는
    width-only 지정으로 원본 비율을 유지하므로 이 값이 이미지를 변형시키지 않는다)."""
    size = image_size(path)
    return (size[0] / size[1]) if size else 1.33


# ---------- 슬라이드 플랜 ----------

def collect_items(blocks):
    """블록들의 항목을 (마커, 텍스트) 쌍으로 평탄화한다. numbered 는 파서가 보존한
    원본 마커를 그대로 쓴다 — 블록이 분절돼도 재번호하지 않아 배지 번호와 일치한다."""
    items = []
    for b in blocks:
        if b["type"] == "bullets":
            items.extend(("•", t) for t in b["items"])
        elif b["type"] == "numbered":
            items.extend((it["marker"], it["text"]) for it in b["items"])
    return items


def sec_visual(sec):
    """절에 시각 요소(캡처 또는 placeholder)가 있는가 — 없으면 병합 후보 개요 절."""
    return any(b["type"] in ("image", "placeholder") for b in sec["blocks"])


LINE_H = 0.275  # 설명 항목 줄당 높이(in) — 11.5pt + space_after 6pt


def _load_badge_positions(img_path):
    """이미지에 대응하는 markers.json 에서 배지별 세로 위치를 읽는다.

    반환: {배지 번호(n): y_frac} — 이미지 전체 높이 대비 배지 상단(0~1).
    markers.json 이 없거나 파싱 실패면 {}. **번호를 키로 돌려주는 이유**: 순서 목록으로
    돌려주면 미발견(found=false) 배지가 섞였을 때 원고 항목과 한 칸씩 밀려 엉뚱한
    밴드로 배분된다. 원고 마커(①=1, '1.'=1)와 번호로 조인해야 안전하다."""
    import json
    stem = os.path.splitext(img_path)[0]
    cands = []
    if stem.endswith("_annotated"):     # resolve_image 는 _annotated 를 우선 반환
        cands.append(stem[:-len("_annotated")] + ".markers.json")
    cands.append(stem + ".markers.json")
    for mp in cands:
        if not os.path.exists(mp):
            continue
        try:
            with open(mp, encoding="utf-8") as f:
                data = json.load(f)
            out = {}
            for m in data.get("markers", []):
                if m.get("found") and m.get("n"):
                    out[int(m["n"])] = float(m.get("y", 0.0))
            return out
        except (OSError, ValueError, KeyError, TypeError):
            continue
    return {}


TILE_FAIL_MSG = {
    "no-pillow": "Pillow 미설치 — 세로 긴 캡처를 분할하지 못해 폭이 줄어든 채 렌더됩니다. "
                 "`pip install Pillow` 후 재생성하세요",
    "no-size": "이미지 치수를 읽지 못했습니다(지원되지 않는 포맷) — PNG 로 다시 저장하세요",
    "error": "이미지 처리 중 오류 — 파일이 손상되었는지 확인하세요",
    "not-needed": "",
    "band-limit": "세로가 과도하게 길어 밴드 상한(6장)에 걸렸습니다 — 밴드가 표준보다 커져 "
                  "폭이 줄어듭니다. 기능 경계로 나눠 개별 캡처(논리 분할)하세요",
}


def warn_tall_tile(img_path, why, sec):
    """세로 긴 캡처의 타일 분할 실패·한계를 stderr 로 알린다 — 조용히 넘어가면
    폭 균일 규격이 무음으로 깨진 산출물이 납품된다."""
    msg = TILE_FAIL_MSG.get(why or "", "")
    if not msg:
        return
    print(f"[build_pptx] 경고: '{sec['num']} {sec['title']}' 의 세로 긴 캡처 "
          f"({os.path.basename(img_path)}) — {msg}", file=sys.stderr)


def _chunk_items(items, width_ea, budget):
    """항목을 슬라이드당 상한(MAX_ITEMS)·줄 예산으로 나눈 청크 목록을 돌려준다.

    일반 컷과 타일 밴드 컷이 **같은 분할 규칙**을 쓰게 하려고 분리했다 — 밴드 경로가
    이 로직을 건너뛰면 항목이 한 장에 몰려 본문이 슬라이드 밖으로 넘친다."""
    chunks = [[]]
    used = 0
    for marker, text in items:
        lines = text_lines(plain(text), width_ea)
        if chunks[-1] and (len(chunks[-1]) >= MAX_ITEMS or used + lines > budget):
            chunks.append([])
            used = 0
        chunks[-1].append((marker, text))
        used += lines

    # 균형 재배분: 그리디 분할(6+1 등)은 마지막 장에 항목 1~2개만 남는 고아
    # 슬라이드를 만들므로, 청크 수는 유지한 채 항목을 균등(4+3 등)하게 다시 나눈다
    if len(chunks) > 1:
        k, n = len(chunks), len(items)
        sizes = [n // k + (1 if i < n % k else 0) for i in range(k)]
        balanced, pos = [], 0
        for size in sizes:
            balanced.append(items[pos:pos + size])
            pos += size
        if all(len(c) <= MAX_ITEMS and
               sum(text_lines(plain(t), width_ea) for _, t in c) <= budget * 1.2
               for c in balanced):
            chunks = balanced
    return chunks


def _marker_no(marker):
    """항목 마커(①·'1.')를 배지 번호(int)로. 알 수 없으면 None."""
    mk = (marker or "").strip()
    if mk and mk[0] in CIRCLED:
        return CIRCLED.index(mk[0]) + 1
    m = re.match(r"^(\d{1,2})\.$", mk)
    return int(m.group(1)) if m else None


def intro_img_y(paras, access):
    """개요·접근 경로를 렌더한 뒤의 이미지 시작 y(in) — render_screen 과 동일 수식.

    개요가 길면 이미지 프레임이 아래로 밀리는데(img_y = max(IMG_Y, y)), 예산을 상수
    IMG_Y 로만 계산하면 그만큼 본문이 슬라이드 밖으로 넘친다."""
    y = 1.2
    if paras or access:
        est = sum(text_lines(plain(p), INTRO_EA) for p in paras)
        y += 0.28 * est + (0.36 if access else 0.02) + 0.08
    return max(IMG_Y.inches, y)


def _merge_small_chunks(chunks, width_ea, budget, last_extra=0):
    """분할 결과에서 인접 컷을 합칠 수 있으면 합친다 — 항목 2~3개짜리 성긴 컷 방지.
    마지막 컷은 ※주의가 함께 렌더되므로 그 몫을 뺀 예산으로 판정한다."""
    if len(chunks) < 2:
        return chunks
    out = [chunks[0]]
    for idx in range(1, len(chunks)):
        ch = chunks[idx]
        is_last = idx == len(chunks) - 1
        cap = max(2, budget - (last_extra if is_last else 0))
        merged = out[-1] + ch
        lines = sum(text_lines(plain(t), width_ea) for _, t in merged)
        if len(merged) <= MAX_ITEMS and lines <= cap:
            out[-1] = merged
        else:
            out.append(ch)
    return out


def _band_budget(frame_h_in, extra_lines=0, img_y_in=None):
    """밴드(또는 상/하 배치) 컷에서 이미지 아래에 남는 설명 줄 예산."""
    top = IMG_Y.inches if img_y_in is None else img_y_in
    avail = TEXT_BOTTOM.inches - top - (0.05 + CAP_H.inches + 0.08)
    return max(2, int((avail - frame_h_in) / LINE_H) - extra_lines)


def _std_frame_h(img_path):
    """상/하 배치에서 이미지 프레임이 차지하는 표준 높이(in)."""
    if img_path and PORTRAIT:
        return min(H_FRAME_W.inches / image_ratio(img_path), PORT_IMG_MAX_H.inches)
    return H_FRAME_H.inches


def _avail_below_intro(img_y):
    """이미지 프레임 상단부터 설명 하한까지, 캡션 몫을 뺀 가용 높이(in)."""
    return TEXT_BOTTOM.inches - img_y - (0.05 + CAP_H.inches + 0.08)


def _seg_layout(visual, img_path, need_lines, img_y=None):
    """세그먼트의 (줄당 문자 수, 설명 줄 예산, 확정 프레임 높이 in|None).

    예산은 정적 상수가 아니라 그 컷의 시각 요소 실높이가 남기는 공간에서 계산한다
    — 이미지가 큰 컷에서 텍스트가 슬라이드 밖으로 넘치는 것을 분할 단계에서 막는다.
    상(이미지)/하(설명) 배치에서 이미지는 축소하지 않고 항상 전폭(표준 프레임)으로
    둔다 — 폭이 매뉴얼 일관성의 앵커이기 때문이다. 설명이 예산을 넘으면 이미지를
    줄이는 대신 분할 단계가 추가 컷으로 나눈다(각 컷이 동일 크기 이미지를 반복).
    세로로 과하게 긴 이미지(비율<TALL_RATIO_MIN)는 이 함수 이전에 타일 분할된다."""
    if img_y is None:
        img_y = IMG_Y.inches
    if visual is None:
        return WIDE_EA, PLAIN_LINES, None
    horizontal = bool(img_path) and image_ratio(img_path) >= 1.45
    if not (PORTRAIT or horizontal):
        # 가로형 좌(이미지)/우(설명) — 우측 컬럼은 이미지 상단부터 하한까지 전부 쓴다
        return SIDE_EA, int((TEXT_BOTTOM.inches - img_y) / LINE_H), None
    # 표준 프레임 높이: 세로형 이미지는 비율 기반 실높이(상한 반영), 그 외는 프레임 규격
    std_h = _std_frame_h(img_path)
    avail = _avail_below_intro(img_y)
    budget = max(2, int((avail - std_h) / LINE_H))
    if not PORTRAIT and need_lines > budget:
        # 가로형 상/하 배치: 이미지는 프레임에 맞춰 넣는 방식(fit)이라 비율에 따라 폭이
        # 이미 달라진다 — 폭 앵커 규격은 세로형 전용이므로, 여기서는 프레임을 하한까지
        # 줄여 한 컷에 담는 편이 낫다(제거하면 컷당 4줄 제한으로 슬라이드가 3배로 늘어난다).
        want_h = avail - need_lines * LINE_H
        if want_h >= IMG_MIN_H:
            return WIDE_EA, need_lines, want_h
        return WIDE_EA, max(2, int((avail - IMG_MIN_H) / LINE_H)), IMG_MIN_H
    return WIDE_EA, budget, std_h


def split_section(sec, draft_dir, shots_dir):
    """절 하나를 1개 이상의 화면 슬라이드 플랜으로 나눈다.

    시각 요소(캡처·placeholder) 경계를 우선한다: 원고 순서대로 각 시각 요소와 그
    뒤에 이어지는 항목들을 한 세그먼트로 묶는다 — 이미지가 여러 장이거나
    placeholder 와 혼재해도 어느 것도 소실되지 않고, 컷-설명 대응이 원고 순서
    그대로 보존된다. 항목 수(MAX_ITEMS)·줄 예산 2차 분할은 세그먼트 안에서만
    일어나며, 분할된 컷들은 그 세그먼트의 시각 요소를 반복 표시한다."""
    blocks = sec["blocks"]
    paras = [b["text"] for b in blocks if b["type"] == "para"]
    access = next((b["text"] for b in blocks if b["type"] == "access"), "")
    notes = [b["text"] for b in blocks if b["type"] == "note"]
    tables = [b for b in blocks if b["type"] == "table"]

    # 시각 요소 경계 세그먼트: (visual_block|None, [(marker, text) ...]) 목록
    segments = []
    cur_visual, cur_items = None, []
    for b in blocks:
        if b["type"] in ("image", "placeholder"):
            if cur_visual is not None or cur_items:
                segments.append((cur_visual, cur_items))
            cur_visual, cur_items = b, []
        elif b["type"] == "bullets":
            cur_items.extend(("•", t) for t in b["items"])
        elif b["type"] == "numbered":
            cur_items.extend((it["marker"], it["text"]) for it in b["items"])
    segments.append((cur_visual, cur_items))

    # 표는 항목과 달리 '실제 도형'이라 줄 예산을 깎아도 자리가 생기지 않는다. 이미지
    # 아래에 표가 들어갈 높이가 안 나오면 표를 앞선 전용 컷으로 뺀다 — 그대로 두면
    # 본문 프레임이 설명 하한 밖으로 밀려 내용이 페이지에서 사라진다. 원고 순서도
    # 개요 → 접근 경로 → 표 → 이미지이므로 표가 앞서는 편이 맞다.
    table_own_cut = False
    if tables and segments[0][0] is not None:
        v0, items0 = segments[0]
        img0 = v0 if v0["type"] == "image" else None
        path0 = resolve_image(img0["src"], draft_dir, shots_dir) if img0 else None
        horiz0 = bool(path0) and image_ratio(path0) >= 1.45
        if PORTRAIT or horiz0:      # 상/하 배치에서만 표가 이미지 아래로 밀린다
            img_y0 = intro_img_y(paras, access)
            tbl_h = tables_height(tables, BODY_W.inches)
            room = _avail_below_intro(img_y0) - _std_frame_h(path0)
            # 표만 겨우 들어가고 설명이 한 줄도 못 실리면 나눈 의미가 없다 — 2줄 여유 요구
            if tbl_h + (2 * LINE_H if items0 else 0) > room:
                table_own_cut = True
                own_room = TEXT_BOTTOM.inches - img_y0
                if tbl_h > own_room:
                    print(f"[build_pptx] 경고: '{sec['num']} {sec['title']}' 표가 한 쪽에 담기지 않습니다"
                          f" (표 {tbl_h:.1f}in > 가용 {own_room:.1f}in) — 행을 줄이거나"
                          " 표를 나눠 주세요", file=sys.stderr)

    plans = []
    last_si = len(segments) - 1
    for si, (visual, items) in enumerate(segments):
        image = visual if visual is not None and visual["type"] == "image" else None
        ph = visual if visual is not None and visual["type"] == "placeholder" else None
        img_path = resolve_image(image["src"], draft_dir, shots_dir) if image else None
        # 표는 첫 컷, ※주의는 마지막 컷에 렌더되므로 그 컷의 필요 줄·예산에 반영한다
        wea = SIDE_EA if (visual is not None and not PORTRAIT
                          and not (bool(img_path) and image_ratio(img_path) >= 1.45)) else WIDE_EA
        # 표는 첫 컷, ※주의는 마지막 컷에만 렌더된다 — 예산도 그 컷에서만 빼야 한다.
        # 모든 컷에서 빼면 앞 컷이 불필요하게 잘게 쪼개진다(밴드 경로와 동일 규칙).
        first_extra = last_extra = 0
        if si == 0 and tables and not table_own_cut:
            first_extra = sum(math.ceil((table_height_est(tb["rows"], BODY_W.inches) + 0.25) / LINE_H)
                              for tb in tables)
        if si == last_si and notes:
            last_extra = sum(text_lines(plain(nt), wea) for nt in notes)
        extra = first_extra + last_extra
        need = sum(text_lines(plain(t), wea) for _, t in items) + extra
        # 개요·접근 경로는 첫 컷에만 실리고 그만큼 이미지가 아래로 밀린다 — 예산도
        # 그 실제 시작 위치에서 계산해야 본문이 슬라이드 밖으로 나가지 않는다
        seg_img_y = intro_img_y(paras, access) if si == 0 else IMG_Y.inches
        width_ea, budget, frame_h = _seg_layout(visual, img_path, need, seg_img_y)

        # 세로 긴 이미지: 폭을 줄이는 대신 표준 비율 밴드로 타일 분할한다(요소 경계
        # 스냅). 폭이 균일성 앵커이므로 모든 밴드가 전폭 6.4in로 렌더된다. 설명은
        # 첫 밴드에 싣고, 이후 밴드는 '(계속 k/N)' 이미지 전용 연속 컷으로 둔다.
        bands = []
        if PORTRAIT and img_path and TALL_RATIO_MIN and image_ratio(img_path) < TALL_RATIO_MIN:
            bands, why = tile_tall_image(img_path, shots_dir)
            if len(bands) < 2 or why == "band-limit":
                warn_tall_tile(img_path, why, sec)
        if len(bands) > 1:
            nb = len(bands)
            base_cap = (image.get("caption") if image else "") or ""
            bh = [(image_size(bp) or (0, 1))[1] or 1 for bp in bands]
            th = sum(bh) or 1
            # 밴드 경계(전체 높이 대비 누적 하단 비율)
            edges, acc = [], 0
            for hpx in bh:
                acc += hpx
                edges.append(acc / th)
            # 배지 y좌표(markers.json)가 있으면 각 배지가 실제로 놓인 밴드에 그 항목을
            # 배분한다 — 원고 마커 번호로 조인하므로 미발견 배지가 섞여도 밀리지 않는다.
            # 좌표가 없거나 하나도 매칭되지 않으면 밴드 높이 비율로 순차 배분(폴백).
            ymap = _load_badge_positions(img_path)
            dist = [[] for _ in range(nb)]
            matched, last_b = 0, 0
            for it in items:
                no = _marker_no(it[0])
                y = ymap.get(no) if no is not None else None
                if y is None:
                    dist[last_b].append(it)      # 매칭 없는 항목은 직전 항목과 같은 밴드에
                    continue
                b = next((bi for bi, bot in enumerate(edges) if y < bot), nb - 1)
                dist[b].append(it)
                last_b = b
                matched += 1
            if matched == 0:
                dist = [[] for _ in range(nb)]
                pos = 0
                for bi in range(nb):
                    cnt = (len(items) - pos) if bi == nb - 1 else round(len(items) * bh[bi] / th)
                    dist[bi] = items[pos:pos + max(0, cnt)]
                    pos += max(0, cnt)

            for bi, bpath in enumerate(bands):
                br = image_ratio(bpath)
                bframe_h = min(BODY_W.inches / br, PORT_IMG_MAX_H.inches)
                # 밴드마다 이미지 높이가 다르므로 예산도 밴드별로 계산하고, 그 예산으로
                # 다시 청크 분할한다 — 배지가 한 밴드에 몰려도 본문이 넘치지 않는다
                # 표는 첫 컷, ※주의는 마지막 컷에만 렌더되므로 **그 컷의 예산에서만** 뺀다
                # — 밴드의 모든 청크에서 빼면 컷이 불필요하게 잘게 쪼개진다
                first_extra = last_extra = 0
                if bi == 0 and si == 0 and tables and not table_own_cut:
                    first_extra = sum(math.ceil((table_height_est(tb["rows"], BODY_W.inches) + 0.25) / LINE_H)
                                      for tb in tables)
                if bi == nb - 1 and si == last_si and notes:
                    last_extra = sum(text_lines(plain(nt), WIDE_EA) for nt in notes)
                # 첫 밴드는 개요·접근 경로 뒤에 오므로 밀린 시작 위치로 예산을 잡는다
                budget_n = _band_budget(bframe_h, 0,
                                        seg_img_y if (bi == 0 and si == 0) else IMG_Y.inches)
                if first_extra and dist[bi]:
                    head = _chunk_items(dist[bi], WIDE_EA, max(2, budget_n - first_extra))[0]
                    rest = dist[bi][len(head):]
                    bchunks = [head] + (_chunk_items(rest, WIDE_EA, budget_n) if rest else [])
                else:
                    bchunks = _chunk_items(dist[bi], WIDE_EA, budget_n)
                if last_extra and bchunks:
                    tail = bchunks[-1]
                    if sum(text_lines(plain(t), WIDE_EA) for _, t in tail) + last_extra > budget_n:
                        bchunks = bchunks[:-1] + _chunk_items(tail, WIDE_EA,
                                                              max(2, budget_n - last_extra))
                bchunks = _merge_small_chunks(bchunks, WIDE_EA, budget_n, last_extra)
                cap = f"{base_cap} ({bi + 1}/{nb})".strip() if base_cap else ""
                for chunk in bchunks:
                    plans.append({
                        "kind": "screen", "sec": sec,
                        "image": image, "img_path": bpath, "ph": None,
                        "horizontal": False,
                        "frame_h": bframe_h,
                        "caption": cap,
                        "items": chunk,
                    })
            continue

        if first_extra and items:
            head = _chunk_items(items, width_ea, max(2, budget - first_extra))[0]
            rest = items[len(head):]
            chunks = [head] + (_chunk_items(rest, width_ea, budget) if rest else [])
        else:
            chunks = _chunk_items(items, width_ea, budget)
        if last_extra and chunks:
            tail = chunks[-1]
            if sum(text_lines(plain(t), width_ea) for _, t in tail) + last_extra > budget:
                chunks = chunks[:-1] + _chunk_items(tail, width_ea, max(2, budget - last_extra))
        chunks = _merge_small_chunks(chunks, width_ea, budget, last_extra)

        for chunk in chunks:
            plans.append({
                "kind": "screen", "sec": sec,
                "image": image, "img_path": img_path, "ph": ph,
                "horizontal": bool(img_path) and image_ratio(img_path) >= 1.45,
                "frame_h": frame_h,
                "items": chunk,
            })

    # 표 전용 선행 컷 — 개요·접근 경로도 여기 실리므로(아래 pi == 0 배치) 원고 순서가
    # 그대로 유지되고, 뒤따르는 이미지 컷은 표준 위치(IMG_Y)에서 시작해 정렬도 맞는다
    if table_own_cut:
        plans.insert(0, {
            "kind": "screen", "sec": sec,
            "image": None, "img_path": None, "ph": None,
            "horizontal": False, "frame_h": None, "items": [],
        })

    # 절 수준 요소 배치: 개요·접근 경로·표는 첫 컷, ※주의는 마지막 컷
    total = len(plans)
    for pi, p in enumerate(plans):
        p["part"] = (pi + 1, total)
        p["paras"] = paras if pi == 0 else []
        p["access"] = access if pi == 0 else ""
        p["tables"] = tables if pi == 0 else []
        p["notes"] = notes if pi == total - 1 else []
    return plans


# 개요 병합 슬라이드 — 렌더(render_sec_stack)와 동일한 수식으로 높이를 추정한다
STACK_HEAD = 0.30      # 절 소제목 줄
STACK_PARA = 0.27      # 개요 문단(12.5pt) 줄당
STACK_ITEM = 0.25      # 항목·주의(11.5/11pt) 줄당
STACK_TABLE_PAD = TABLE_PAD
STACK_GAP = 0.18       # 절 사이 간격
COMBINE_Y = 1.25       # 병합 본문 시작 y (in) — 가용 높이(COMBINE_BUDGET)는 방향 프로파일이 정한다

# table_height_est 는 draft_parser 로 이관 — 원고 린터가 pptx 의존 없이 같은 수식을 쓴다


def sec_stack_height(sec):
    """개요 절이 병합 슬라이드에서 차지할 높이(in) 추정."""
    h = STACK_HEAD
    for b in sec["blocks"]:
        if b["type"] == "para":
            h += text_lines(plain(b["text"]), INTRO_EA) * STACK_PARA
        elif b["type"] == "access":
            h += STACK_PARA
        elif b["type"] == "note":
            h += text_lines(plain(b["text"]), WIDE_EA) * STACK_ITEM
        elif b["type"] == "table":
            h += table_height_est(b["rows"], BODY_W.inches) + STACK_TABLE_PAD
        elif b["type"] == "bullets":
            h += sum(text_lines(plain(t), WIDE_EA) for t in b["items"]) * STACK_ITEM
        elif b["type"] == "numbered":
            h += sum(text_lines(plain(it["text"]), WIDE_EA) for it in b["items"]) * STACK_ITEM
    return h


def combine_overview_runs(sections, draft_dir, shots_dir, ch):
    """장 안의 절들을 순서대로 플랜으로 바꾸되, 연속된 개요 절(시각 요소 없음)은
    분량이 예산 안이면 한 슬라이드로 병합한다. 절 경계에서만 나눈다.

    제목만 있고 본문 블록이 없는 절(3단 번호의 부모/그룹 제목)은 슬라이드를 만들지
    않는다 — 헤더만 있는 빈 장표가 생기기 때문이다. 대신 CONTENTS 번호를 다음 실제
    슬라이드(보통 첫 하위 절)에 위임한다(also_secs)."""
    plans, run, pending = [], [], []

    def attach(plan_item):
        # 제목만 있는 절들의 목차 번호를 이 슬라이드로 위임한다
        if pending:
            plan_item.setdefault("also_secs", []).extend(pending)
            pending.clear()

    def flush():
        if not run:
            return
        groups, cur, used = [], [], 0.0
        for sec in run:
            h = sec_stack_height(sec) + (STACK_GAP if cur else 0)
            if cur and used + h > COMBINE_BUDGET:
                groups.append(cur)
                cur, used = [], 0.0
                h = sec_stack_height(sec)
            cur.append(sec)
            used += h
        groups.append(cur)
        combined = [g for g in groups if len(g) > 1]
        gi = 0
        for group in groups:
            # 홀로 남은 절(분량 초과 포함)은 기존 단독 슬라이드 경로가 레이아웃을 보장한다
            if len(group) == 1:
                new = split_section(group[0], draft_dir, shots_dir)
                if new:
                    attach(new[0])
                plans.extend(new)
            else:
                gi += 1
                item = {"kind": "combined", "ch": ch, "secs": group,
                        "part": (gi, len(combined))}
                attach(item)
                plans.append(item)
        run.clear()

    for sec in sections:
        if not sec["blocks"]:
            flush()
            pending.append(sec)
        elif sec_visual(sec):
            flush()
            new = split_section(sec, draft_dir, shots_dir)
            if new:
                attach(new[0])
            plans.extend(new)
        else:
            run.append(sec)
    flush()
    if pending and plans:
        # 장 끝에 남은 제목-only 절 — 마지막 슬라이드에 위임한다
        plans[-1].setdefault("also_secs", []).extend(pending)
    return plans


def build_plan(doc, draft_dir, shots_dir):
    screens = []
    for ch in doc["chapters"]:
        screens.append({"kind": "divider", "ch": ch})
        screens.extend(combine_overview_runs(ch["sections"], draft_dir, shots_dir, ch))

    toc_items = []
    for ch in doc["chapters"]:
        toc_items.append(("ch", f"{ch['num']}. {ch['title']}", ch))
        for sec in ch["sections"]:
            toc_items.append(("sec", f"{sec['num']} {sec['title']}", sec))
    per_col, cols = TOC_PER_COL, TOC_COLS
    toc_pages = max(1, math.ceil(len(toc_items) / (per_col * cols)))

    plan = [{"kind": "cover"}] + [{"kind": "contents", "page": i} for i in range(toc_pages)] + screens

    slide_no = {}
    for idx, item in enumerate(plan, start=1):
        if item["kind"] == "divider":
            slide_no.setdefault(f"ch:{item['ch']['num']}", idx)
        elif item["kind"] == "screen" and item["part"][0] == 1:
            slide_no.setdefault(f"sec:{item['sec']['num']}", idx)
        elif item["kind"] == "combined":
            for sec in item["secs"]:
                slide_no.setdefault(f"sec:{sec['num']}", idx)
        for sec in item.get("also_secs", []):
            # 제목만 있는 부모 절 — 위임받은 슬라이드 번호를 목차에 표기한다
            slide_no.setdefault(f"sec:{sec['num']}", idx)
    return plan, toc_items, per_col, cols, slide_no


# ---------- 렌더 ----------

def render_cover(prs, doc, args):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, DARK)
    audience, version, date = parse_meta(doc["meta"])
    audience = args.audience or audience or "사용자용"
    version = args.version or version
    date = args.date or date
    title = args.title or doc["title"] or "사용자 매뉴얼"
    label = args.cover_label

    # 제목이 "관리자 매뉴얼 — 시스템명" 꼴이면 프로젝트명만 남긴다 — 표지 규격은
    # [매뉴얼 구분(분리형일 때만)] + [프로젝트명] 이고, 구분 표기는 label 이 담당한다
    m = re.match(r"^(.{0,20}매뉴얼)\s*[—\-–:]\s*(.+)$", title)
    if m:
        title = m.group(2).strip()

    # 세로 위치는 슬라이드 높이 비율로 잡는다 — 가로/세로 어느 방향에서도 같은 균형
    vh = lambda frac: Inches(SLIDE_H.inches * frac)
    cover_w = SLIDE_W - Inches(1.6)
    if label:
        # 독자 영역별 매뉴얼이 분리된 세트 — 1행 "〈독자〉 매뉴얼", 개행 후 프로젝트명
        tf = add_text(slide, Inches(0.8), vh(0.245), cover_w, Inches(0.5))
        add_para(tf, label, 20, color=DARK_SOFT, align=PP_ALIGN.CENTER)
    tf = add_text(slide, Inches(0.8), vh(0.32), cover_w, Inches(1.2), anchor=MSO_ANCHOR.MIDDLE)
    add_para(tf, title, 34, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    if not label:
        # 단일 매뉴얼 — 프로젝트명만 크게, 아래에 대상 부제
        tf = add_text(slide, Inches(0.8), vh(0.493), cover_w, Inches(0.6))
        add_para(tf, f"{audience} 사용자 매뉴얼" if "매뉴얼" not in audience else audience,
                 16, color=DARK_SOFT, align=PP_ALIGN.CENTER)
    add_rect(slide, 0, vh(0.653), SLIDE_W, Pt(3), ACCENT_ON_DARK)  # 어두운 표지 위 — 원색 유지
    meta_line = " · ".join(v for v in (audience, f"버전 {version}" if version else "", date) if v)
    tf = add_text(slide, Inches(0.8), vh(0.707), cover_w, Inches(0.5))
    add_para(tf, meta_line, 12, color=DARK_SOFTER, align=PP_ALIGN.CENTER)


def render_contents(prs, page_idx, toc_items, per_col, cols, slide_no):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    tf = add_text(slide, Inches(0.7), Inches(0.5), Inches(6), Inches(0.7))
    add_para(tf, "CONTENTS", 26, color=DARK, bold=True)
    add_rect(slide, Inches(0.72), Inches(1.25), Inches(1.6), Pt(2.5), ACCENT)

    start = page_idx * per_col * cols
    chunk = toc_items[start:start + per_col * cols]
    col_w = TOC_COL_W
    # 마지막(오버플로) 페이지에서 좌측 컬럼에만 몰리지 않도록 좌우 균등 분배
    half = max(1, math.ceil(len(chunk) / cols))
    for c in range(cols):
        col_items = chunk[c * half:(c + 1) * half]
        if not col_items:
            break
        tf = add_text(slide, Inches(0.72) + c * (col_w + Inches(0.3)), Inches(1.6), col_w,
                      SLIDE_H - Inches(2.1))
        for kind, label, ref in col_items:
            key = f"ch:{ref['num']}" if kind == "ch" else f"sec:{ref['num']}"
            no = slide_no.get(key, "")
            if kind == "ch":
                add_para(tf, [(f"{label}", True), (f"   ····  {no}", False)], 14,
                         color=DARK, space_after=6)
            else:
                add_para(tf, [(f"{label}", False), (f"   ····  {no}", False)], 11.5,
                         color=TEXT, space_after=5)


def render_divider(prs, ch):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, DARK)
    vh = lambda frac: Inches(SLIDE_H.inches * frac)
    tf = add_text(slide, Inches(1.0), vh(0.36), SLIDE_W - Inches(2.0), Inches(1.4))
    add_para(tf, f"{ch['num']}.", 52, color=ACCENT_ON_DARK, bold=True)  # 어두운 간지 위
    tf = add_text(slide, Inches(1.0), vh(0.533), SLIDE_W - Inches(2.0), Inches(1.0))
    add_para(tf, ch["title"], 30, color=WHITE, bold=True)


def render_header(slide, ch, sec, part, right_label=True):
    add_rect(slide, 0, 0, SLIDE_W, Inches(1.0), DARK)
    suffix = f" ({part[0]}/{part[1]})" if part[1] > 1 else ""
    # 세로형은 폭이 좁아 제목 폭을 우측 장 라벨과 겹치지 않게 줄인다
    title_w = SLIDE_W - Inches(3.4) if PORTRAIT else Inches(9.2)
    tf = add_text(slide, Inches(0.55), Inches(0.18), title_w, Inches(0.7), anchor=MSO_ANCHOR.MIDDLE)
    p = tf.paragraphs[0]
    set_para(p, [(f"{sec['num']}  ", False), (sec["title"] + suffix, True)], 20, color=WHITE)
    p.runs[0].font.color.rgb = HEAD_NUM
    if right_label:
        label_w = Inches(2.35) if PORTRAIT else Inches(3.3)
        tf = add_text(slide, SLIDE_W - label_w - Inches(0.43), Inches(0.33), label_w, Inches(0.4))
        add_para(tf, f"{ch['num']}. {ch['title']}", 11, color=DARK_SOFT, align=PP_ALIGN.RIGHT)


def render_page_no(slide, no):
    tf = add_text(slide, SLIDE_W - Inches(0.78), SLIDE_H - Inches(0.45), Inches(0.6), Inches(0.35))
    add_para(tf, str(no), 10, color=MUTED, align=PP_ALIGN.RIGHT)


def apply_picture_border(pic):
    """그림 서식: 검은색 실선 테두리(두께 기본값) — 흰 배경 캡처와 슬라이드의 경계를 살린다.
    그림자는 쓰지 않는다(뷰어별 렌더 편차·템플릿 경로 누락 문제로 테두리로 일원화)."""
    pic.line.color.rgb = RGBColor(0x00, 0x00, 0x00)


def render_items(tf, items, size):
    for marker, text in items:
        segs = [(f"{marker} ", False)] + parse_inline(text)
        add_para(tf, segs, size, space_after=6)


def render_sec_stack(slide, sec, y):
    """개요 절 하나를 병합 슬라이드의 y 위치에 그리고 다음 y 를 돌려준다.
    높이 수식은 sec_stack_height 와 반드시 일치해야 한다(겹침 방지의 근거)."""
    blocks = sec["blocks"]
    paras = [b["text"] for b in blocks if b["type"] == "para"]
    access = next((b["text"] for b in blocks if b["type"] == "access"), "")
    notes = [b["text"] for b in blocks if b["type"] == "note"]
    tables = [b for b in blocks if b["type"] == "table"]
    items = collect_items(blocks)

    head_h = STACK_HEAD + sum(text_lines(plain(p), INTRO_EA) for p in paras) * STACK_PARA \
        + (STACK_PARA if access else 0)
    tf = add_text(slide, BODY_X, Inches(y), BODY_W, Inches(head_h))
    add_para(tf, [(f"{sec['num']}  ", True), (sec["title"], True)], 15, color=DARK, space_after=5)
    for para_text in paras:
        add_para(tf, parse_inline(para_text), 12.5)
    if access:
        add_para(tf, [("접근 경로  ", True), (access, False)], 11.5, color=ACCENT)
    y += head_h

    for tb in tables:
        render_table(slide, tb["rows"], BODY_X, Inches(y), BODY_W)
        y += table_height_est(tb["rows"], BODY_W.inches) + STACK_TABLE_PAD

    if items or notes:
        body_h = sum(text_lines(plain(t), WIDE_EA) for _, t in items) * STACK_ITEM \
            + sum(text_lines(plain(n), WIDE_EA) for n in notes) * STACK_ITEM
        tf = add_text(slide, BODY_X, Inches(y), BODY_W, Inches(max(0.3, body_h)))
        render_items(tf, items, 11.5)
        for note in notes:
            add_para(tf, parse_inline(note), 11, color=NOTE, space_after=4)
        y += body_h
    return y


def render_combined(prs, item, page_no):
    """이미지 없는 연속 개요 절 묶음을 한 슬라이드로 — 절 제목을 소제목으로 스택한다."""
    ch = item["ch"]
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    render_header(slide, ch, {"num": ch["num"], "title": ch["title"]}, item["part"],
                  right_label=False)
    render_page_no(slide, page_no)
    y = COMBINE_Y
    for sec in item["secs"]:
        y = render_sec_stack(slide, sec, y) + STACK_GAP


def render_table(slide, rows, x, y, w):
    n_rows, n_cols = len(rows), max(len(r) for r in rows)
    height = Inches(0.35) * n_rows
    shape = slide.shapes.add_table(n_rows, n_cols, x, y, w, height)
    table = shape.table
    for ri, row in enumerate(rows):
        for ci in range(n_cols):
            cell = table.cell(ri, ci)
            cell.text = row[ci] if ci < len(row) else ""
            for p in cell.text_frame.paragraphs:
                for r in p.runs:
                    _set_font(r, 11, bold=(ri == 0), color=DARK if ri == 0 else TEXT)
    return shape


def render_screen(prs, plan_item, ch_of, page_no):
    sec = plan_item["sec"]
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    render_header(slide, ch_of[sec["num"]], sec, plan_item["part"])
    render_page_no(slide, page_no)

    # 개요 문단과 접근 경로를 한 프레임에 넣는다 — 줄 수 추정이 빗나가도
    # 프레임 안에서 이어지므로 서로 겹칠 수 없다
    y = Inches(1.2)
    if plan_item["paras"] or plan_item["access"]:
        est = sum(text_lines(plain(p), INTRO_EA) for p in plan_item["paras"])
        tf = add_text(slide, BODY_X, y, BODY_W,
                      Inches(0.3) * (est + (1 if plan_item["access"] else 0)))
        for para_text in plan_item["paras"]:
            add_para(tf, parse_inline(para_text), 12.5)
        if plan_item["access"]:
            add_para(tf, [("접근 경로  ", True), (plan_item["access"], False)], 11.5, color=ACCENT)
        y += Inches(0.28) * est + (Inches(0.36) if plan_item["access"] else Inches(0.02)) + Inches(0.08)

    image, ph = plan_item["image"], plan_item["ph"]
    img_path, horizontal = plan_item["img_path"], plan_item["horizontal"]

    # 이미지 프레임 상단은 전 장표 공통(IMG_Y)으로 고정한다 — 개요 길이에 따라 캡처
    # 크기·위치가 장표마다 달라 보이는 것을 막기 위함이다. 개요가 예약 공간(2줄+접근
    # 경로)을 초과하는 예외에서만 겹침 방지를 위해 프레임을 아래로 민다.
    img_y = max(IMG_Y, y)
    if img_y > IMG_Y:
        print(f"[build_pptx] 경고: '{sec['num']} {sec['title']}' 개요가 표준 예약(2줄)을 넘어 "
              "이미지 프레임이 아래로 밀렸습니다 — 개요를 1~2문장으로 줄이면 전 장표 정렬이 유지됩니다",
              file=sys.stderr)

    if image or ph:
        # 상(이미지)/하(설명) 배치 여부 — 세로형은 항상, 가로형은 가로 비율 캡처만
        top_layout = PORTRAIT or (horizontal and img_path)
        frame_w, frame_h = (H_FRAME_W, H_FRAME_H) if top_layout else (V_FRAME_W, V_FRAME_H)
        # 분할 단계가 확정한 프레임 높이(설명 수용 위한 축소 포함)를 그대로 쓴다 —
        # 분할 예산과 렌더 기하가 같은 수식을 공유해야 겹침·넘침이 없다
        if plan_item.get("frame_h"):
            frame_h = Inches(plan_item["frame_h"])
        frame_x = (SLIDE_W - frame_w) / 2 if top_layout else BODY_X

        if img_path and image_size(img_path) is None:
            # 치수를 모르는 포맷 — 폭만 지정해 렌더러가 원본 비율을 유지하게 하고,
            # 실측 높이가 상한을 넘으면 비율 유지로 축소한다 (비율 가정 변형 금지)
            pic = slide.shapes.add_picture(img_path, frame_x, img_y, width=frame_w)
            if pic.height > frame_h:
                pic.width = int(pic.width * (frame_h / pic.height))
                pic.height = int(frame_h)
                pic.left = int(frame_x + (frame_w - pic.width) / 2)
            frame_h = pic.height  # 캡션·설명 시작 위치가 실측 높이를 따르도록
            apply_picture_border(pic)
        elif img_path:
            ratio = image_ratio(img_path)
            if PORTRAIT:
                # 세로형: 본문 폭을 가득 채우고 높이는 비율 유지 — 뷰포트를 통일한
                # 캡처라면 전 장표에서 동일 크기가 된다. 프레임 상한(비율 실높이
                # 또는 설명 수용 축소값)을 넘으면 비율 유지로 줄인다.
                w = frame_w
                h = w / ratio
                if h > frame_h:
                    h = frame_h
                    w = h * ratio
                frame_h = h  # 캡션·설명 시작 위치가 실제 이미지 높이를 따르도록
            else:
                # 가로형: 비율 유지로 프레임에 맞추고(fit) 프레임 중앙에 배치 — 같은
                # 비율의 캡처는 어느 장표에서든 동일한 크기로 렌더된다
                w = min(frame_w, frame_h * ratio)
                h = w / ratio
            px = frame_x + (frame_w - w) / 2
            py = img_y + (frame_h - h) / 2
            apply_picture_border(slide.shapes.add_picture(img_path, px, py, width=w, height=h))
        else:
            # placeholder 이거나, 이미지 블록은 있으나 파일이 누락된 경우 — 어느 쪽이든
            # 프레임과 동일한 크기의 회색 안내 상자로 렌더한다
            info = ph or {"scr": (image or {}).get("scr", ""),
                          "name": f"이미지 파일 누락: {(image or {}).get('src', '')}"}
            box = add_rect(slide, frame_x, img_y, frame_w, frame_h, PH_BG,
                           line=RGBColor(0xC9, 0xCE, 0xD6))
            btf = box.text_frame
            btf.word_wrap = True
            btf.vertical_anchor = MSO_ANCHOR.MIDDLE
            set_para(btf.paragraphs[0], "화면 이미지 추후 삽입", 14, color=PH_TX,
                     bold=True, align=PP_ALIGN.CENTER)
            set_para(btf.add_paragraph(), f"{info['scr']}  {info['name']}", 11, color=PH_TX,
                     align=PP_ALIGN.CENTER)

        cap_y = img_y + frame_h + Inches(0.05)  # 캡션도 프레임 하단 고정 위치
        # 타일 밴드는 plan 이 캡션을 직접 지정(원 캡션 + "(계속 k/N)")한다
        caption = plan_item.get("caption")
        if caption is None:
            caption = (image.get("caption") if image else "") or (f"[사진 -] {ph['name']} (추후 삽입)" if ph else "")
        if caption:
            tf = add_text(slide, frame_x, cap_y, frame_w, CAP_H)
            add_para(tf, caption, 9.5, color=MUTED, align=PP_ALIGN.CENTER)

        if top_layout:
            text_x, text_y, text_w = BODY_X, cap_y + CAP_H + Inches(0.08), BODY_W
        else:
            text_x, text_y, text_w = SIDE_X, img_y, SIDE_W
    else:
        text_x, text_y, text_w = BODY_X, y, BODY_W

    # 표를 먼저 그리고, 설명 텍스트 프레임은 표 아래에서 시작한다 (겹침 방지)
    ty = text_y
    for tb in plan_item["tables"]:
        render_table(slide, tb["rows"], text_x, ty, text_w)
        ty += Inches(table_height_est(tb["rows"], text_w.inches)) + Inches(0.25)
    tf = add_text(slide, text_x, ty, text_w, max(Inches(0.4), TEXT_BOTTOM - ty))
    render_items(tf, plan_item["items"], 11.5)
    for note in plan_item["notes"]:
        add_para(tf, parse_inline(note), 11, color=NOTE, space_after=4)


# ---------- 자체 검증 ----------

def self_check(prs, combined_idx=frozenset()):
    problems = []
    pic_widths = []
    for idx, slide in enumerate(prs.slides, start=1):
        item_count = 0
        for shape in slide.shapes:
            # 스크린샷 렌더 폭 수집 — 매뉴얼 전체에서 캡처가 균일 폭으로 들어가야
            # 일관성이 유지된다(폭이 균일성 앵커). shape_type 13 = PICTURE.
            if getattr(shape, "shape_type", None) == 13 and shape.width:
                pic_widths.append(shape.width)
            if not shape.has_text_frame:
                continue
            for p in shape.text_frame.paragraphs:
                text = "".join(r.text for r in p.runs)
                # "[버튼](부연)" 은 원고의 정상 표기이므로 이미지 문법(![)과 볼드(**)만 검사한다
                if "**" in text or "![" in text:
                    problems.append(f"슬라이드 {idx}: 마크다운 잔재 의심 — {text[:40]}")
                stripped = text.strip()
                first = stripped[:1]
                # 장 번호는 zero-padded("03. ")로 렌더되므로 0으로 시작하는 번호는 제외
                if first == "•" or (first and first in CIRCLED) or re.match(r"^(?!0)\d{1,2}\.\s", stripped):
                    item_count += 1
        # 항목 밀도 규칙(6개 분할)은 화면 슬라이드 대상 — 개요 병합 슬라이드는 여러 절의
        # 짧은 항목이 합산되므로 제외한다 (높이는 병합 예산이 이미 보장)
        if item_count > MAX_ITEMS and idx not in combined_idx:
            problems.append(f"슬라이드 {idx}: 설명 항목 {item_count}개 (분할 기준 {MAX_ITEMS} 초과)")
        # 텍스트 프레임 이탈 — 표·개요가 길어 본문이 밀리면 내용이 페이지 밖으로 사라진다.
        # 산출물 검증(verify_pptx)과 같은 규칙을 빌드 시점에도 두어 즉시 드러나게 한다.
        for shape in slide.shapes:
            if (getattr(shape, "has_text_frame", False) and shape.top is not None
                    and shape.height is not None and shape.top + shape.height > SLIDE_H):
                over = (shape.top + shape.height - SLIDE_H) / 914400
                problems.append(f"슬라이드 {idx}: 텍스트 프레임이 페이지 아래로 {over:.2f}in "
                                "넘칩니다 — 표가 길거나 개요가 예약 공간을 초과했습니다")
                break

    # 렌더 폭 균일 게이트 — 세로형 전용(가로형은 좌/우·상/하 배치라 이미지 폭이
    # 본래 다르다). 세로 긴 캡처가 조용히 폭 축소되어 다른 장표와 크기가 달라지는
    # 것을 잡는다. 폭이 매뉴얼 일관성의 앵커다.
    if PORTRAIT and pic_widths:
        lo, hi = min(pic_widths), max(pic_widths)
        if hi and (hi - lo) / hi > 0.03:
            problems.append(f"세로형 스크린샷 렌더 폭 편차 {round((hi - lo) / hi * 100)}% "
                            f"({round(lo / 914400, 2)}~{round(hi / 914400, 2)}in) — 세로 긴 캡처의 "
                            "폭 축소 의심. 표준 뷰포트로 재캡처하거나 논리 분할 권장")
        elif lo < BODY_W * 0.90:
            problems.append(f"세로형 스크린샷 폭 {round(lo / 914400, 2)}in < 전폭 90% "
                            "— 폭 축소 렌더(타일 실패 또는 비표준 비율)")
    return problems


def main():
    ap = argparse.ArgumentParser(description="manual-draft.md → 매뉴얼 PPTX")
    ap.add_argument("--draft", required=True, help="manual-draft.md 경로")
    ap.add_argument("--screenshots", help="스크린샷 디렉토리 (기본: draft 위치의 screenshots/)")
    ap.add_argument("--out", required=True, help="산출 pptx 경로")
    ap.add_argument("--title", help="표지 제목 (기본: 원고 # 제목)")
    ap.add_argument("--audience", help='표지 대상 표기 (예: "관리자용")')
    ap.add_argument("--version", help="표지 버전 표기")
    ap.add_argument("--date", help="표지 날짜 표기")
    ap.add_argument("--orientation", choices=["landscape", "portrait"], default="portrait",
                    help="슬라이드 방향: portrait=A4 세로(기본) / landscape=16:9 가로")
    ap.add_argument("--theme", choices=sorted(THEMES), default="navy",
                    help="색 테마: navy=네이비+블루(기본) / forest=딥그린+틸 / charcoal=차콜+오렌지")
    ap.add_argument("--theme-from", metavar="TEMPLATE_PPTX",
                    help="참고 템플릿 pptx 에서 색 테마(dark·accent)를 추출해 적용 — "
                         "레이아웃·규격은 번들 그대로(서식 복제 아님). 추출 실패 시 --theme 폴백")
    ap.add_argument("--cover-label",
                    help='표지 1행 매뉴얼 구분 표기(예: "관리자 매뉴얼") — 독자 영역별로 '
                         "매뉴얼이 분리된 세트일 때만 지정. 지정 시 표지가 [구분]+[프로젝트명] "
                         "2행이 되고, 미지정 시 프로젝트명만 표기한다")
    ap.add_argument("--skip-validate", action="store_true", help="원고 사전 검증을 건너뛴다")
    args = ap.parse_args()

    apply_orientation(args.orientation == "portrait")
    apply_theme(args.theme)
    if args.theme_from:
        got = theme_from_template(args.theme_from)
        if got:
            print(f"[build_pptx] 참고 템플릿 색 테마 적용: dark=#{got[0]} accent=#{got[1]} "
                  f"({os.path.basename(args.theme_from)}) — 레이아웃·규격은 번들 표준")
        else:
            print(f"[build_pptx] 경고: 템플릿 테마 추출 실패({args.theme_from}) — "
                  f"--theme {args.theme} 로 진행", file=sys.stderr)

    if not os.path.exists(args.draft):
        fail(f"원고 없음: {args.draft}")
    draft_dir = os.path.dirname(os.path.abspath(args.draft))
    shots_dir = args.screenshots or os.path.join(draft_dir, "screenshots")

    doc = parse_draft(args.draft)
    if not doc["chapters"]:
        fail("장(## NN. 제목)을 찾지 못했습니다 — manual-template.md 규약을 확인하세요")

    # 원고 사전 검증 게이트 — 잘못된 입력이 결정론적 빌더를 통과해
    # 잘못된 구조로 산출되는 것을 생성 전에 막는다
    if not args.skip_validate:
        import validate_draft
        with open(args.draft, encoding="utf-8-sig") as f:
            raw = f.read()
        errors, warns = validate_draft.validate(doc, draft_dir, shots_dir, raw_text=raw)
        for w in warns:
            print(f"[build_pptx] 원고 WARN: {w}")
        if errors:
            for e in errors:
                print(f"[build_pptx] 원고 ERROR: {e}", file=sys.stderr)
            fail(f"원고 검증 실패({len(errors)}건) — 원고를 수정하거나 --skip-validate 로 우회하세요")

    ch_of = {}
    for ch in doc["chapters"]:
        for sec in ch["sections"]:
            ch_of[sec["num"]] = ch

    plan, toc_items, per_col, cols, slide_no = build_plan(doc, draft_dir, shots_dir)

    prs = Presentation()
    prs.slide_width, prs.slide_height = SLIDE_W, SLIDE_H
    for idx, item in enumerate(plan, start=1):
        if item["kind"] == "cover":
            render_cover(prs, doc, args)
        elif item["kind"] == "contents":
            render_contents(prs, item["page"], toc_items, per_col, cols, slide_no)
            render_page_no(prs.slides[-1], idx)
        elif item["kind"] == "divider":
            render_divider(prs, item["ch"])
        elif item["kind"] == "combined":
            render_combined(prs, item, idx)
        else:
            render_screen(prs, item, ch_of, idx)

    combined_idx = {i for i, it in enumerate(plan, start=1) if it["kind"] == "combined"}
    problems = self_check(prs, combined_idx)
    os.makedirs(os.path.dirname(os.path.abspath(args.out)), exist_ok=True)
    prs.save(args.out)

    missing = len({id(it["ph"]) for it in plan if it["kind"] == "screen" and it["ph"]})
    print(f"[build_pptx] 저장 완료: {args.out} (슬라이드 {len(plan)}장, placeholder {missing}건)")

    # 이미지 소실 안전망 — 원고의 모든 image 블록이 최소 한 컷에 배치됐는지 대조한다
    rendered = {id(it["image"]) for it in plan if it["kind"] == "screen" and it["image"]}
    lost = [b["src"] for ch in doc["chapters"] for sec in ch["sections"]
            for b in sec["blocks"] if b["type"] == "image" and id(b) not in rendered]
    if lost:
        print(f"[build_pptx] 경고: 원고의 이미지 {len(lost)}건이 어떤 슬라이드에도 "
              f"배치되지 않았습니다: {lost}", file=sys.stderr)

    # 배지·테두리 파이프라인 미사용 감지 — 세션이 옛 지침으로 돌아도 조립 시점에 잡는다
    if os.path.isdir(shots_dir):
        originals = [f for f in os.listdir(shots_dir)
                     if f.lower().endswith(".png") and "_annotated" not in f]
        annotated = [f for f in os.listdir(shots_dir) if "_annotated" in f]
        if originals and not annotated:
            print("[build_pptx] 경고: screenshots에 _annotated 파일이 0건 — 배지·강조 테두리 "
                  "파이프라인(cdp_capture.py --mark → annotate_screenshot.py)이 사용되지 않았습니다. "
                  "skill의 표준 형식(번호 배지+테두리)이 빠진 산출물입니다.", file=sys.stderr)
    if problems:
        print("[build_pptx] 자체 검증 경고:")
        for pr in problems:
            print(f"  - {pr}")
    else:
        print("[build_pptx] 자체 검증 통과 (마크다운 잔재·항목 초과 없음)")


if __name__ == "__main__":
    main()
