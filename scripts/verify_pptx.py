"""산출 pptx 규격 검증기 — 생성 경로와 무관하게 최종 파일을 기계 검증한다.

번들 빌더(build_pptx.py)는 self_check 를 내장하지만, 참고 템플릿 서식 복제(전용
skill/자작 빌더) 경로는 규격(output-formats.md)을 우회해도 잡는 장치가 없었다 —
실제로 자작 템플릿 빌더 산출물에서 테두리 0건·캡션 누락·비율 강제 fit 이 무경고로
납품될 뻔한 사례가 있다. 이 검증기는 **pptx 파일 자체**를 열어 규격을 검사하므로
어떤 생성기로 만들었든 동일 게이트를 통과해야 한다.

검사 항목:
  ERROR — 마크다운 잔재(**, ![), 스크린샷 검은 테두리 누락, 이미지 비율 왜곡(>3%),
          (--draft) 원고 이미지 소실
  WARN  — 세로형 렌더 폭 편차(>3%), 슬라이드당 항목 6개 초과, 그림 있는 슬라이드의
          [사진 N] 캡션 부재, 미확정 마킹 잔존, (--draft) placeholder 잔존 불일치

사용 예:
  python verify_pptx.py 관리자매뉴얼_시스템명_20260724.pptx \
      [--draft manual-work/manual-draft.md] [--screenshots manual-work/screenshots]

종료 코드: 0 통과(경고 포함) / 1 ERROR 존재 / 2 파일 없음·python-pptx 미설치
"""

import argparse
import os
import re
import sys

if hasattr(sys.stdout, "reconfigure"):
    sys.stdout.reconfigure(encoding="utf-8")
    sys.stderr.reconfigure(encoding="utf-8")

sys.path.insert(0, os.path.dirname(os.path.abspath(__file__)))
from draft_parser import parse_draft, resolve_image, CIRCLED

try:
    from pptx import Presentation
    from pptx.oxml.ns import qn
except ImportError:
    print("[verify_pptx] python-pptx 미설치: pip install python-pptx", file=sys.stderr)
    sys.exit(2)

MAX_ITEMS = 6
UNRESOLVED_RE = re.compile(r"확인 필요|확정 전|TBD|미정")
PHOTO_RE = re.compile(r"\[사진\s*[\d-]+\]")


def has_border(pic):
    """그림에 실선 테두리(a:ln + solidFill)가 지정되어 있는가."""
    sp = pic._element.spPr
    ln = sp.find(qn("a:ln"))
    if ln is None:
        return False
    return ln.find(qn("a:solidFill")) is not None


def native_ratio(pic):
    """그림 원본(임베드 blob)의 가로/세로 비율. 알 수 없으면 None."""
    try:
        w, h = pic.image.size  # px
        return (w / h) if h else None
    except Exception:
        return None


def verify(path, draft=None, shots_dir=None):
    prs = Presentation(path)
    W, H = prs.slide_width, prs.slide_height
    portrait = H > W
    errors, warns = [], []

    pic_widths = []
    n_pics = 0
    bordered = 0
    for idx, slide in enumerate(prs.slides, start=1):
        item_count = 0
        slide_has_pic = False
        slide_has_caption = False
        for shape in slide.shapes:
            if getattr(shape, "shape_type", None) == 13:  # PICTURE
                n_pics += 1
                slide_has_pic = True
                pic_widths.append(shape.width)
                if has_border(shape):
                    bordered += 1
                else:
                    errors.append(f"슬라이드 {idx}: 스크린샷 검은 테두리 누락 — 규격은 전 캡처 "
                                  "검은 실선(output-formats.md)")
                nr = native_ratio(shape)
                if nr and shape.height:
                    rr = shape.width / shape.height
                    if abs(rr - nr) / nr > 0.03:
                        errors.append(f"슬라이드 {idx}: 이미지 비율 왜곡 {nr:.2f}→{rr:.2f} "
                                      "— 폭·높이 강제 지정(비율 유지 위반)")
            if not shape.has_text_frame:
                continue
            for p in shape.text_frame.paragraphs:
                text = "".join(r.text for r in p.runs)
                if "**" in text or "![" in text:
                    errors.append(f"슬라이드 {idx}: 마크다운 잔재 — {text[:40]}")
                if PHOTO_RE.search(text):
                    slide_has_caption = True
                stripped = text.strip()
                first = stripped[:1]
                if first == "•" or (first and first in CIRCLED) or re.match(r"^(?!0)\d{1,2}\.\s", stripped):
                    item_count += 1
        if item_count > MAX_ITEMS:
            warns.append(f"슬라이드 {idx}: 설명 항목 {item_count}개 (분할 기준 {MAX_ITEMS} 초과)")
        if slide_has_pic and not slide_has_caption:
            warns.append(f"슬라이드 {idx}: 그림은 있는데 [사진 N] 캡션이 없습니다")

    # 세로형 렌더 폭 균일(폭 = 일관성 앵커)
    if portrait and pic_widths:
        lo, hi = min(pic_widths), max(pic_widths)
        if hi and (hi - lo) / hi > 0.03:
            warns.append(f"세로형 스크린샷 렌더 폭 편차 {round((hi - lo) / hi * 100)}% "
                         f"({round(lo / 914400, 2)}~{round(hi / 914400, 2)}in) — 균일 폭 규격 위반 의심")

    # 미확정 마킹
    all_text = "\n".join("".join(r.text for r in p.runs)
                         for s in prs.slides for sh in s.shapes if sh.has_text_frame
                         for p in sh.text_frame.paragraphs)
    hits = sorted(set(UNRESOLVED_RE.findall(all_text)))
    if hits:
        warns.append(f"미확정 마킹 잔존: {hits}")

    # 원고 대조(선택): 이미지 소실·placeholder 잔존
    if draft and os.path.exists(draft):
        doc = parse_draft(draft)
        ddir = os.path.dirname(os.path.abspath(draft))
        shots = shots_dir or os.path.join(ddir, "screenshots")
        want = 0
        for ch in doc["chapters"]:
            for sec in ch["sections"]:
                for b in sec["blocks"]:
                    if b["type"] == "image" and resolve_image(b["src"], ddir, shots):
                        want += 1
        # 타일 밴드는 이미지가 늘어나므로 '부족'만 오류로 본다
        if n_pics < want:
            errors.append(f"원고 이미지 {want}건 중 pptx 에 {n_pics}건만 배치 — 이미지 소실 의심")
        n_ph_draft = sum(1 for ch in doc["chapters"] for sec in ch["sections"]
                         for b in sec["blocks"] if b["type"] == "placeholder")
        n_ph_pptx = all_text.count("추후 삽입")
        if n_ph_draft and not n_ph_pptx:
            warns.append(f"원고 placeholder {n_ph_draft}건이 pptx 에 표기되지 않음 — 무음 누락 의심")

    size_lab = f"{W / 914400:.2f}x{H / 914400:.2f}in ({'세로' if portrait else '가로'})"
    return errors, warns, {"slides": len(prs.slides._sldIdLst), "pics": n_pics,
                           "bordered": bordered, "size": size_lab}


def main():
    ap = argparse.ArgumentParser(description="산출 pptx 규격 검증 (생성 경로 무관)")
    ap.add_argument("pptx", help="검증할 pptx 경로")
    ap.add_argument("--draft", help="원고 경로(이미지 소실·placeholder 대조)")
    ap.add_argument("--screenshots", help="스크린샷 디렉토리 (기본: draft 위치의 screenshots/)")
    args = ap.parse_args()

    if not os.path.exists(args.pptx):
        print(f"[verify_pptx] 파일 없음: {args.pptx}", file=sys.stderr)
        sys.exit(2)

    errors, warns, info = verify(args.pptx, args.draft, args.screenshots)
    print(f"[verify_pptx] {os.path.basename(args.pptx)} — {info['size']}, "
          f"슬라이드 {info['slides']}장, 그림 {info['pics']}개(테두리 {info['bordered']})")
    for w in warns:
        print(f"[verify_pptx] WARN: {w}")
    for e in errors:
        print(f"[verify_pptx] ERROR: {e}", file=sys.stderr)
    print(f"[verify_pptx] 결과: ERROR {len(errors)}건 / WARN {len(warns)}건")
    sys.exit(1 if errors else 0)


if __name__ == "__main__":
    main()
