# -*- coding: utf-8 -*-
"""표 기하 보정 1단계 — 측정용 pptx 를 만든다.

빌더가 실제로 쓰는 render_table 을 그대로 호출하므로, 측정 대상은 산출물과 동일한
조판 조건(11pt, 본문 폭, 셀 여백)에 놓인다. 각 슬라이드에 표 하나를 두고 본문 행은
모두 같은 텍스트로 채운다 — 그 조건의 행 높이를 여러 번 재기 위해서다.

측정 조합
  한글  : 열 수(2·3·4) x 글자 수 19단계 — 줄바꿈 경계를 촘촘히 훑는다
  ASCII : 소문자·대문자·넓은 글자(m·w)·좁은 글자(i·l·j)·숫자·한영 혼용
          x 열 수(2·3) x 글자 수 6단계 — 글자 폭 가중을 검증한다

사용:
  python tools/make_table_probe.py --out-dir <작업폴더>
  → <작업폴더>/probe.pptx, <작업폴더>/probe_meta.json

다음 단계는 tools/measure_table.ps1 (README.md 참고).
"""

import argparse
import io
import json
import os
import sys

if hasattr(sys.stdout, "reconfigure"):
    sys.stdout.reconfigure(encoding="utf-8")

sys.path.insert(0, os.path.join(os.path.dirname(os.path.abspath(__file__)), "..", "scripts"))
import build_pptx as B                                    # noqa: E402
from pptx import Presentation                             # noqa: E402
from pptx.util import Inches                              # noqa: E402

# 글자 폭 가중을 검증하기 위한 표본 — 클래스별로 순수하게 채운다
ASCII_SAMPLES = {
    "en_low":  "abcdefghijklmnopqrstuvwxyz" * 4,
    "en_up":   "ABCDEFGHIJKLMNOPQRSTUVWXYZ" * 4,
    "en_wide": "mmmmmmmmmmwwwwwwwwww" * 5,
    "en_narw": "iiiiiiiiiilllllllllljjjj" * 4,
    "digit":   "0123456789" * 10,
    "mix":     "발송 status 확인 OK, 2026-07-24 처리 완료 " * 4,
}
HANGUL_BASE = "가나다라마바사아자차카타파하" * 10
HANGUL_NCH = (2, 4, 6, 8, 10, 12, 14, 16, 18, 20, 24, 28, 32, 36, 40, 48, 56, 64, 80)
ASCII_NCH = (10, 20, 30, 45, 60, 90)


def build(out_dir):
    B.apply_orientation(True)                             # 세로형(A4) 기준으로 보정한다
    prs = Presentation()
    prs.slide_width, prs.slide_height = B.SLIDE_W, B.SLIDE_H
    meta = []

    def add(kind, ncol, text, n_body):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        rows = [[f"머리{c + 1}" for c in range(ncol)]] + [[text] * ncol for _ in range(n_body)]
        B.render_table(slide, rows, Inches(0.55), Inches(1.0), B.BODY_W)
        meta.append({"kind": kind, "ncol": ncol, "nch": len(text), "text": text,
                     "nrows": len(rows)})

    for ncol in (2, 3, 4):
        for nch in HANGUL_NCH:
            add("hangul", ncol, HANGUL_BASE[:nch], 4)
    for kind, base in ASCII_SAMPLES.items():
        for ncol in (2, 3):
            for nch in ASCII_NCH:
                add(kind, ncol, base[:nch], 3)

    os.makedirs(out_dir, exist_ok=True)
    pptx_path = os.path.join(out_dir, "probe.pptx")
    meta_path = os.path.join(out_dir, "probe_meta.json")
    prs.save(pptx_path)
    with open(meta_path, "w", encoding="utf-8") as f:
        json.dump(meta, f, ensure_ascii=False)
    print(f"[probe] 표 {len(meta)}개 생성 → {pptx_path}")
    print(f"[probe] 다음: powershell -File tools/measure_table.ps1 "
          f"-Path \"{os.path.abspath(pptx_path)}\" -Out \"{os.path.join(os.path.abspath(out_dir), 'measured.csv')}\"")


def main():
    ap = argparse.ArgumentParser(description="표 기하 보정용 측정 pptx 생성")
    ap.add_argument("--out-dir", required=True, help="probe.pptx·probe_meta.json 을 둘 폴더")
    build(ap.parse_args().out_dir)


if __name__ == "__main__":
    main()
